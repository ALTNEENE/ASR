from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parent
MEDIA = ROOT / "presentation_assets" / "odt_media"
OUT = ROOT / "ASR_Diabetes_Monitoring_Presentation.pptx"


W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(13, 31, 45)
TEAL = RGBColor(24, 132, 117)
MINT = RGBColor(225, 246, 240)
GREEN = RGBColor(47, 127, 122)
ORANGE = RGBColor(235, 149, 50)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(95, 111, 125)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(212, 221, 228)


def media(name: str) -> Path:
    path = MEDIA / name
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def blank(prs: Presentation):
    return prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def add_title(slide, title, subtitle=None, section=None, dark=False):
    color = WHITE if dark else INK
    accent = ORANGE if dark else TEAL
    if section:
        add_text(slide, section.upper(), Inches(0.62), Inches(0.34), Inches(5.5), Inches(0.26), 9, accent, True)
    add_text(slide, title, Inches(0.6), Inches(0.58), Inches(9.5), Inches(0.58), 26, color, True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.62), Inches(1.12), Inches(9.6), Inches(0.36), 11, MINT if dark else MUTED)
    add_rect(slide, Inches(0.62), Inches(1.55), Inches(1.1), Inches(0.04), accent)


def add_footer(slide, idx, section="A.S.R Diabetes Monitoring"):
    add_text(slide, section, Inches(0.62), Inches(7.08), Inches(5.5), Inches(0.22), 8, MUTED)
    add_text(slide, f"{idx:02d}", Inches(12.05), Inches(7.08), Inches(0.65), Inches(0.22), 8, MUTED, True, PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=14, color=INK, spacing=0.94):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = color
        p.space_after = Pt(6 * spacing)
        p._p.get_or_add_pPr().set("marL", "228600")
        p._p.get_or_add_pPr().set("indent", "-137160")
    return box


def add_numbered(slide, items, x, y, w, h, size=13):
    for i, item in enumerate(items, 1):
        y_pos = y + Inches(0.5 * (i - 1))
        add_rect(slide, x, y_pos + Inches(0.02), Inches(0.28), Inches(0.28), TEAL, radius=True)
        add_text(slide, str(i), x + Inches(0.02), y_pos + Inches(0.05), Inches(0.24), Inches(0.18), 8, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, item, x + Inches(0.43), y_pos, w - Inches(0.43), Inches(0.42), size, INK)


def add_card(slide, title, body, x, y, w, h, accent=TEAL, title_size=13, body_size=10):
    add_rect(slide, x, y, w, h, WHITE, LINE, radius=True)
    add_rect(slide, x, y, Inches(0.07), h, accent)
    add_text(slide, title, x + Inches(0.22), y + Inches(0.18), w - Inches(0.36), Inches(0.28), title_size, INK, True)
    if body:
        add_text(slide, body, x + Inches(0.22), y + Inches(0.54), w - Inches(0.36), h - Inches(0.65), body_size, MUTED)


def add_image_fit(slide, path, x, y, w, h, border=True):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    dx = x + int((w - dw) / 2)
    dy = y + int((h - dh) / 2)
    if border:
        add_rect(slide, x, y, w, h, WHITE, LINE, radius=True)
    return slide.shapes.add_picture(str(path), dx, dy, width=dw, height=dh)


def add_image_tile(slide, path, caption, x, y, w, h, cap_size=8):
    add_image_fit(slide, path, x, y, w, h - Inches(0.28), True)
    add_text(slide, caption, x, y + h - Inches(0.22), w, Inches(0.18), cap_size, MUTED, True, PP_ALIGN.CENTER)


def add_table(slide, rows, x, y, w, h, font_size=8):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    for col in range(len(rows[0])):
        table.columns[col].width = int(w / len(rows[0]))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL if r == 0 else (RGBColor(239, 247, 245) if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Aptos"
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else INK
    return table_shape


def add_background(slide, dark=False):
    add_rect(slide, 0, 0, W, H, NAVY if dark else LIGHT)
    if not dark:
        add_rect(slide, 0, 0, Inches(0.16), H, TEAL)


def make_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides = []

    # 1
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s, True)
    add_image_fit(s, media("image10.jpg"), Inches(8.0), Inches(0.75), Inches(4.75), Inches(2.65), True)
    add_image_fit(s, media("image1.png"), Inches(0.78), Inches(0.65), Inches(1.24), Inches(1.24), False)
    add_text(s, "Remote Diabetes Monitoring and Follow-up System", Inches(0.78), Inches(2.05), Inches(7.0), Inches(1.0), 31, WHITE, True)
    add_text(s, "A.S.R smart web platform for proactive diabetes monitoring, awareness, and report sharing", Inches(0.82), Inches(3.18), Inches(6.8), Inches(0.56), 14, MINT)
    add_rect(s, Inches(0.82), Inches(4.05), Inches(2.1), Inches(0.08), ORANGE)
    add_text(s, "Karary University | Faculty of Computer Science and Information Technology", Inches(0.82), Inches(4.45), Inches(7.0), Inches(0.35), 11, MINT)
    add_text(s, "Prepared by: Ahmed Tajelsir, Raghad Osman, Saja Heraika", Inches(0.82), Inches(4.86), Inches(7.4), Inches(0.35), 10, MINT)
    add_text(s, "Supervisor: Dr. Nizar Ali Fideel", Inches(0.82), Inches(5.2), Inches(5.0), Inches(0.3), 10, MINT)
    add_text(s, "Graduation Project Presentation", Inches(9.0), Inches(6.55), Inches(3.2), Inches(0.35), 11, ORANGE, True, PP_ALIGN.RIGHT)
    add_footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Introduction and Motivation", "Why remote diabetes monitoring matters", "Context")
    add_card(s, "Diabetes is a long-term follow-up problem", "Patients need continuous tracking, feedback, and education to avoid complications.", Inches(0.72), Inches(1.95), Inches(3.75), Inches(1.15), TEAL)
    add_card(s, "Traditional follow-up has access barriers", "Remote and resource-constrained areas can face delays, specialist shortages, and high travel costs.", Inches(4.82), Inches(1.95), Inches(3.75), Inches(1.15), ORANGE)
    add_card(s, "Digital platforms can close the distance", "Web monitoring, dashboards, alerts, reports, and AI awareness tools create faster patient-clinician communication.", Inches(8.92), Inches(1.95), Inches(3.75), Inches(1.15), GREEN)
    add_bullets(s, [
        "The A.S.R platform shifts diabetes care from passive recording to proactive monitoring and prevention.",
        "The system targets confirmed diabetic patients, preventive users, and awareness-only visitors through different workflows.",
        "Core idea: connect patient-entered readings with automated classification, personalized guidance, and shareable medical reports."
    ], Inches(0.9), Inches(3.65), Inches(7.6), Inches(2.0), 15)
    add_image_fit(s, media("image25.png"), Inches(9.0), Inches(3.45), Inches(3.35), Inches(2.2), True)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Research Problem", "Challenges addressed by the proposed system", "Problem")
    add_bullets(s, [
        "Limited access to preventive health services and continuous follow-up, especially in developing countries such as Sudan.",
        "Existing digital tools often monitor readings without integrated decision support or practical reporting.",
        "Previous studies focus heavily on confirmed diabetic patients and leave preventive/awareness users underserved.",
        "Manual follow-up increases unnecessary visits and makes longitudinal patient data harder for doctors to review."
    ], Inches(0.85), Inches(1.85), Inches(6.6), Inches(3.7), 16)
    add_card(s, "Project response", "A web-based A.S.R ecosystem that records readings, classifies status instantly, routes users by risk profile, provides AI-supported awareness, and generates reports for remote clinical follow-up.", Inches(8.0), Inches(2.05), Inches(4.2), Inches(2.45), TEAL, 15, 12)
    add_card(s, "Main affected users", "Diabetic patients, people at risk or unsure of their status, general awareness visitors, and admins managing content.", Inches(8.0), Inches(4.85), Inches(4.2), Inches(1.25), ORANGE, 13, 11)
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Objectives, Scope, and Boundaries", "What the project is designed to deliver", "Scope")
    add_numbered(s, [
        "Design a preventive website for regular tracking of diabetes health indicators.",
        "Support early detection through automated classification of abnormal readings.",
        "Improve access to follow-up services for remote or underserved patients.",
        "Reduce unnecessary hospital visits through remote reports and communication.",
        "Raise awareness and support self-management with educational content and AI guidance.",
        "Provide a simple interface suitable for users with limited technical experience."
    ], Inches(0.8), Inches(1.82), Inches(6.5), Inches(3.3), 11)
    add_card(s, "Scope", "Remote diabetes monitoring and follow-up web system connecting patients, reports, awareness content, and smart guidance.", Inches(7.65), Inches(1.86), Inches(4.6), Inches(1.15), TEAL)
    add_card(s, "Time boundary", "Project period: 2025-2026.", Inches(7.65), Inches(3.25), Inches(2.15), Inches(0.95), ORANGE)
    add_card(s, "Domain boundary", "Medical equipment and diabetes follow-up workflows. The system supports education and monitoring, not medical diagnosis replacement.", Inches(10.1), Inches(3.25), Inches(2.15), Inches(0.95), GREEN)
    add_card(s, "System boundary", "User-entered readings are analyzed by stored thresholds and smart assistant prompts; direct IoT sensor integration is a future recommendation.", Inches(7.65), Inches(4.45), Inches(4.6), Inches(1.35), TEAL)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Research Methodology", "Descriptive analytical method with practical implementation", "Method")
    add_bullets(s, [
        "Descriptive analysis: identifies healthcare challenges, user needs, and current diabetes awareness/follow-up practices.",
        "Analytical review: compares earlier remote monitoring studies to extract strengths, weaknesses, and adoption barriers.",
        "Data collection: interviews with health professionals plus scientific sources, reports, and previous studies.",
        "System proposal and implementation: converts findings into workflows, database design, APIs, dashboards, and reports."
    ], Inches(0.8), Inches(1.82), Inches(6.4), Inches(3.8), 15)
    add_card(s, "Use case scenario", "Login -> questionnaire -> reading entry and context -> smart analysis -> results and recommendations -> PDF report.", Inches(7.65), Inches(1.9), Inches(4.6), Inches(1.3), TEAL, 13, 11)
    add_card(s, "User branching", "Diabetic path for clinical tools, preventive path for unsure/at-risk users, and awareness path for education.", Inches(7.65), Inches(3.45), Inches(4.6), Inches(1.3), ORANGE, 13, 11)
    add_card(s, "Analysis output", "Instant classification: below normal, normal, simple high, or sharp rise needing attention.", Inches(7.65), Inches(5.0), Inches(4.6), Inches(0.95), GREEN, 13, 11)
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Literature Review Highlights", "Digital health patterns found across previous work", "Literature")
    add_card(s, "Remote Patient Monitoring", "Digital tools support glucose and wound follow-up, enabling periodic review without constant hospital visits.", Inches(0.8), Inches(1.8), Inches(3.65), Inches(1.35), TEAL)
    add_card(s, "Technology Integration", "EHR, CDSS, IoT, blockchain, and ML can improve access, prediction, security, and clinical decisions.", Inches(4.85), Inches(1.8), Inches(3.65), Inches(1.35), ORANGE)
    add_card(s, "Patient Engagement", "User-friendly platforms improve adherence, satisfaction, and self-management behavior.", Inches(8.9), Inches(1.8), Inches(3.65), Inches(1.35), GREEN)
    add_card(s, "Infrastructure Challenges", "Weak internet, smartphone availability, connectivity issues, and low-tech deployment remain adoption barriers.", Inches(0.8), Inches(3.65), Inches(3.65), Inches(1.35), TEAL)
    add_card(s, "Evidence Gaps", "Many studies lack long-term follow-up, control groups, quantitative comparison, or practical deployment.", Inches(4.85), Inches(3.65), Inches(3.65), Inches(1.35), ORANGE)
    add_card(s, "Web Platform Value", "Online diabetes management can improve glycemic control and increase access in underserved settings.", Inches(8.9), Inches(3.65), Inches(3.65), Inches(1.35), GREEN)
    add_footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Related Work Comparison", "What earlier systems achieved and where they were limited", "Related Work")
    rows = [
        ["Study", "Approach", "Key result", "Limitation"],
        ["Istepanian 2009", "Mobile/Bluetooth glucose monitoring", "RCT design; early mobile monitoring", "Connectivity issues; dropout; short follow-up"],
        ["Ahmed 2017", "EHR + CDSS architecture", "Better data access and decision support design", "Interoperability and infrastructure challenges"],
        ["Keegan 2023", "Patient-centered wound monitoring", "84% adherence; 94% satisfaction", "Small sample and short duration"],
        ["Ratta 2024", "Blockchain + IoT + ML ecosystem", "Security, accuracy, and predictive insights", "Complex deployment in low-tech areas"],
        ["Sawyer 2025", "Five-stage RPM program", "HbA1c improved from 10.4% to 7.0%", "Requires clinical resources and structured program"],
        ["Gardner 2025", "Digital health review/feasibility", "Wound area reduction and high satisfaction", "Limited practical evaluation/statistical comparison"],
    ]
    add_table(s, rows, Inches(0.72), Inches(1.75), Inches(11.9), Inches(4.8), 7)
    add_footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Research Gap and A.S.R Contribution", "How the proposed platform differs", "Gap")
    add_card(s, "Data integration gap", "A.S.R links reading analysis, database records, AI guidance, and report generation in one workflow.", Inches(0.8), Inches(1.85), Inches(5.45), Inches(1.15), TEAL)
    add_card(s, "Reliability and usability gap", "The platform uses a web-based interface, validation, and simple navigation to avoid device-specific friction.", Inches(6.8), Inches(1.85), Inches(5.45), Inches(1.15), ORANGE)
    add_card(s, "Preventive vision gap", "Unlike many systems focused only on diabetic patients, A.S.R provides preventive and awareness paths.", Inches(0.8), Inches(3.35), Inches(5.45), Inches(1.15), GREEN)
    add_card(s, "Practical evaluation gap", "The research produced a functional implementation with real database schema, APIs, dashboard, PDF output, and share links.", Inches(6.8), Inches(3.35), Inches(5.45), Inches(1.15), TEAL)
    add_card(s, "Clinical bridge", "The output is designed to organize patient readings into a doctor-readable report, supporting remote consultation.", Inches(3.0), Inches(5.0), Inches(7.25), Inches(0.95), ORANGE, 14, 11)
    add_footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Functional Requirements", "Main services the system must provide", "Requirements")
    left = [
        "User registration, login, logout, remember session, and Google OAuth sign-in.",
        "Smart questionnaire to classify the user into diabetic, preventive, or awareness path.",
        "Patient profile capture: age, gender, weight, diagnosis date, diabetes type, and treatment data.",
        "Blood glucose reading entry with unit, timing, meal context, notes, symptoms, medication, and weight.",
        "Automated glucose classification by age/context thresholds: low, normal, or high.",
    ]
    right = [
        "Dashboard with latest reading, totals, trends, health status, and personalized advice.",
        "AI assistant for diabetes awareness, nutrition, symptoms, and general guidance.",
        "Medication search and user medication management.",
        "Reading history with filtering, printing, PDF export, email sending, and shareable report links.",
        "Admin workflow for managing users and health knowledge/content."
    ]
    add_bullets(s, left, Inches(0.85), Inches(1.75), Inches(5.7), Inches(4.95), 12)
    add_bullets(s, right, Inches(6.85), Inches(1.75), Inches(5.7), Inches(4.95), 12)
    add_footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Non-Functional Requirements", "Quality attributes reflected in the implementation", "Requirements")
    cards = [
        ("Usability", "Clear Arabic-first screens, simple forms, and guided path routing.", TEAL),
        ("Accessibility", "Browser-based access for remote patients and low-friction use on common devices.", ORANGE),
        ("Reliability", "Deterministic glucose evaluator plus validation before storing readings.", GREEN),
        ("Security and privacy", "Sessions, password hashing, OAuth, environment-based secrets, expiring share tokens.", TEAL),
        ("Performance", "Lightweight PHP APIs, indexed MySQL tables, and focused dashboard requests.", ORANGE),
        ("Maintainability", "Separated config, APIs, reusable evaluator class, migrations, and Composer dependencies.", GREEN),
        ("Data integrity", "Unique users, profile constraints, normalized medications, readings, tips, and report shares.", TEAL),
        ("Extensibility", "Ready for IoT meters, ML prediction, emergency alerts, and stronger encryption.", ORANGE),
    ]
    for i, (title, body, accent) in enumerate(cards):
        col = i % 2
        row = i // 2
        add_card(s, title, body, Inches(0.85 + col * 6.1), Inches(1.75 + row * 1.22), Inches(5.55), Inches(0.92), accent, 12, 9)
    add_footer(s, 10)

    # 11
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Tools and Technology Stack", "Implementation stack confirmed from the project files", "Stack")
    add_card(s, "Frontend", "HTML, CSS, JavaScript, responsive Arabic UI, autocomplete and theme scripts.", Inches(0.85), Inches(1.8), Inches(3.65), Inches(1.3), TEAL)
    add_card(s, "Backend", "PHP 8+, sessions, modular APIs, Composer autoloading, reusable glucose evaluator.", Inches(4.85), Inches(1.8), Inches(3.65), Inches(1.3), ORANGE)
    add_card(s, "Database", "MySQL/MariaDB tables for users, profiles, patient data, readings, medications, tips, AI evaluations, and report shares.", Inches(8.85), Inches(1.8), Inches(3.65), Inches(1.3), GREEN)
    add_card(s, "AI and integrations", "Gemini/Groq smart assistant APIs, Google OAuth 2.0 authentication, share links, email report delivery.", Inches(0.85), Inches(3.45), Inches(3.65), Inches(1.3), TEAL)
    add_card(s, "Reporting", "Dompdf for PDF reports and PHPMailer for sending reports by email.", Inches(4.85), Inches(3.45), Inches(3.65), Inches(1.3), ORANGE)
    add_card(s, "Development tools", "VS Code, StarUML, XAMPP/Apache, Composer, SQL migrations, Git/Vercel deployment files.", Inches(8.85), Inches(3.45), Inches(3.65), Inches(1.3), GREEN)
    add_footer(s, 11)

    # 12
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "System Architecture Diagram", "Layered web platform, data tier, and external APIs", "Diagrams")
    add_image_fit(s, media("image2.png"), Inches(0.8), Inches(1.75), Inches(7.2), Inches(4.05), True)
    add_bullets(s, [
        "Presentation layer: browser-based A.S.R interface.",
        "Application layer: authentication, questionnaire routing, reading evaluator, smart integration, and report generator.",
        "Data tier: users/patients, glucose readings, knowledge/tips, medications, AI evaluations, and share tokens.",
        "External APIs: AI assistant APIs and Google OAuth 2.0."
    ], Inches(8.35), Inches(1.9), Inches(3.95), Inches(3.4), 12)
    add_footer(s, 12)

    # 13
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Use Case and Activity Diagrams", "Actors, workflows, and branching paths", "Diagrams")
    add_image_tile(s, media("image3.png"), "Figure 3.2 - Use Case Diagram", Inches(0.8), Inches(1.75), Inches(5.75), Inches(4.65))
    add_image_tile(s, media("image4.png"), "Figure 3.3 - Activity Diagram", Inches(6.95), Inches(1.75), Inches(5.35), Inches(4.65))
    add_footer(s, 13)

    # 14
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Class Diagram", "Core data objects and relationships", "Diagrams")
    add_image_fit(s, media("image5.png"), Inches(0.9), Inches(1.65), Inches(5.25), Inches(4.95), True)
    add_bullets(s, [
        "User and Patient store authentication and diabetes-specific profile information.",
        "GlucoseReading records values, context, psychological state, medication, and classification output.",
        "SmartAssistant supports analysis and health guidance.",
        "Report produces PDF/shareable summaries for remote doctor review.",
        "Admin manages user accounts and system knowledge content."
    ], Inches(6.65), Inches(1.85), Inches(5.6), Inches(3.9), 13)
    add_footer(s, 14)

    # 15
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Administration and Signup Sequences", "Account creation and admin management flows", "Diagrams")
    add_image_tile(s, media("image6.png"), "Figure 3.5 - User Signup", Inches(0.62), Inches(1.7), Inches(4.05), Inches(4.8), 7)
    add_image_tile(s, media("image7.png"), "Figure 3.6 - Manage User", Inches(4.77), Inches(1.7), Inches(3.8), Inches(4.8), 7)
    add_image_tile(s, media("image8.png"), "Figure 3.7 - Manage System Content", Inches(8.68), Inches(1.7), Inches(3.8), Inches(4.8), 7)
    add_footer(s, 15)

    # 16
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Smart Glucose Evaluation Sequence", "Reading submission, AI suggestions, and saved output", "Diagrams")
    add_image_fit(s, media("image9.png"), Inches(0.8), Inches(1.75), Inches(6.0), Inches(3.95), True)
    add_bullets(s, [
        "Patient submits reading value plus measurement context through the dashboard.",
        "Evaluator retrieves age/type target ranges from the database.",
        "Smart assistant returns personalized recommendations.",
        "System stores reading, classification, and generated guidance.",
        "Final output is shown to the patient and can be included in reports."
    ], Inches(7.25), Inches(1.9), Inches(4.95), Inches(3.7), 12)
    add_footer(s, 16)

    # 17
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "UI Walkthrough: Entry, Intake, and Awareness", "The first user journey through the platform", "UI")
    add_image_tile(s, media("image10.jpg"), "Welcome / Entry Point", Inches(0.62), Inches(1.7), Inches(2.95), Inches(2.05), 7)
    add_image_tile(s, media("image11.png"), "Questionnaire / Intake", Inches(3.78), Inches(1.7), Inches(2.95), Inches(2.05), 7)
    add_image_tile(s, media("image12.png"), "Smart Assistant", Inches(6.94), Inches(1.7), Inches(2.95), Inches(2.05), 7)
    add_image_tile(s, media("image13.png"), "Awareness Resources", Inches(10.1), Inches(1.7), Inches(2.55), Inches(2.05), 7)
    add_bullets(s, [
        "Users start from a clear welcome page and move into a questionnaire-based routing step.",
        "Confirmed diabetic users continue to the clinical monitoring path.",
        "Non-diabetic or unsure users receive awareness resources and guided assistant support."
    ], Inches(1.05), Inches(4.42), Inches(11.2), Inches(1.15), 13)
    add_footer(s, 17)

    # 18
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "UI Walkthrough: Authentication", "Registration and login for the diabetic path", "UI")
    add_image_tile(s, media("image14.png"), "Registration Interface", Inches(0.95), Inches(1.75), Inches(5.55), Inches(3.55), 8)
    add_image_tile(s, media("image15.png"), "Login Interface", Inches(6.9), Inches(1.75), Inches(5.55), Inches(3.55), 8)
    add_bullets(s, [
        "Phone/password authentication is supported for direct access.",
        "Google OAuth provides a faster single sign-on path.",
        "After login, routing depends on survey completion and diabetes status."
    ], Inches(1.3), Inches(5.65), Inches(10.5), Inches(0.7), 12)
    add_footer(s, 18)

    # 19
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "UI Walkthrough: Dashboard and Personalized Advice", "Monitoring summary after sign-in", "UI")
    add_image_tile(s, media("image16.png"), "Integrated Patient Dashboard", Inches(0.95), Inches(1.75), Inches(5.6), Inches(3.75), 8)
    add_image_tile(s, media("image17.png"), "Personalized Advice", Inches(6.95), Inches(1.75), Inches(4.55), Inches(3.75), 8)
    add_bullets(s, [
        "Dashboard summarizes profile status, warning signs, recent readings, and quick actions.",
        "Advice adapts to age, gender, diabetes information, and measurement history."
    ], Inches(1.05), Inches(5.78), Inches(11.0), Inches(0.55), 12)
    add_footer(s, 19)

    # 20
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "UI Walkthrough: Clinical Tools and Preventive Path", "Recording, reporting, sharing, and unsure-user guidance", "UI")
    add_image_tile(s, media("image18.png"), "Glucose Logging Tool", Inches(0.62), Inches(1.7), Inches(2.65), Inches(3.0), 7)
    add_image_tile(s, media("image19.png"), "Evaluation PDF Report", Inches(3.42), Inches(1.7), Inches(2.95), Inches(3.0), 7)
    add_image_tile(s, media("image20.png"), "Readings History", Inches(6.52), Inches(1.7), Inches(2.55), Inches(3.0), 7)
    add_image_tile(s, media("image21.png"), "Guided Tools", Inches(9.2), Inches(1.7), Inches(1.55), Inches(3.0), 7)
    add_image_tile(s, media("image22.png"), "Unsure User Guidance", Inches(10.95), Inches(1.7), Inches(1.65), Inches(3.0), 7)
    add_bullets(s, [
        "Patients enter measured glucose values with meal timing, notes, medication, symptoms, and weight.",
        "The system stores readings, classifies status, supports filtering, and generates doctor-ready outputs.",
        "Preventive users receive short intake questions and assistant guidance instead of clinical monitoring tools."
    ], Inches(0.95), Inches(5.25), Inches(11.4), Inches(1.0), 11)
    add_footer(s, 20)

    # 21
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s)
    add_title(s, "Results and Recommendations", "Implemented outcomes and future roadmap", "Results")
    add_image_tile(s, media("image23.png"), "Comparison Chart", Inches(0.62), Inches(1.72), Inches(3.65), Inches(2.35), 7)
    add_image_tile(s, media("image24.png"), "Workflow Distribution", Inches(4.55), Inches(1.72), Inches(3.25), Inches(2.35), 7)
    add_image_tile(s, media("image25.png"), "Predictive Analytics Roadmap", Inches(8.05), Inches(1.72), Inches(4.55), Inches(2.35), 7)
    add_card(s, "Results", "The six research objectives were implemented as working features: monitoring, early detection, remote access, PDF reports, AI awareness, and usable interfaces.", Inches(0.8), Inches(4.45), Inches(5.6), Inches(1.15), TEAL, 12, 9)
    add_card(s, "Recommendations", "Add direct smart glucose-meter synchronization, ML-based predictive alerts, GPS emergency escalation, and stronger end-to-end medical data privacy.", Inches(6.75), Inches(4.45), Inches(5.6), Inches(1.15), ORANGE, 12, 9)
    add_footer(s, 21)

    # 22
    s = prs.slides.add_slide(blank(prs)); slides.append(s); add_background(s, True)
    add_image_fit(s, media("image1.png"), Inches(0.82), Inches(0.65), Inches(1.2), Inches(1.2), False)
    add_text(s, "Thank You", Inches(0.82), Inches(1.95), Inches(4.6), Inches(0.75), 38, WHITE, True)
    add_text(s, "Questions and Discussion", Inches(0.86), Inches(2.72), Inches(4.6), Inches(0.35), 16, MINT)
    add_rect(s, Inches(0.86), Inches(3.25), Inches(1.55), Inches(0.08), ORANGE)
    refs = [
        "IDF. (2021). IDF Diabetes Atlas, 10th ed.",
        "WHO. (2021). Global Strategy on Digital Health 2020-2025.",
        "WHO. (2023). Diabetes Fact Sheet.",
        "Hou et al. (2021). Web-based diabetes self-management interventions.",
        "Istepanian et al. (2009). Mobile phone telemonitoring for glycemic control.",
        "Ahmed et al. (2017). RPM system engineering for diabetes management.",
        "Keegan et al. (2023). Patient-centered remote wound monitoring.",
        "Ratta et al. (2024). Blockchain-ML ecosystem for IoT diabetes monitoring.",
        "Sawyer et al. (2025). Successful RPM program for diabetes.",
        "Gardner et al. (2025). Digital health technology in Asia-Pacific diabetes management.",
        "ADA. (2022). Standards of Medical Care in Diabetes.",
        "Greenwood et al. (2017). Digital health interventions for diabetes management."
    ]
    add_text(s, "Selected References", Inches(6.25), Inches(0.72), Inches(5.8), Inches(0.42), 18, WHITE, True)
    add_bullets(s, refs, Inches(6.25), Inches(1.25), Inches(5.9), Inches(4.95), 8, MINT, 0.45)
    add_text(s, "A.S.R Remote Diabetes Monitoring and Follow-up System", Inches(0.86), Inches(5.88), Inches(5.1), Inches(0.34), 12, MINT, True)
    add_footer(s, 22)

    prs.save(OUT)
    return OUT, len(slides)


if __name__ == "__main__":
    output, count = make_presentation()
    print(f"Generated {output}")
    print(f"Slides: {count}")
